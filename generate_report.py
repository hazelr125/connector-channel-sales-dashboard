import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

df = pd.read_csv('connector_leads.csv')

#Funnel metrics
def funnel_summary(df, groupby_col):
    g = df.groupby(groupby_col).agg(
        leads=('lead_id', 'count'),
        approvals=('approved', 'sum'),
        disbursals=('disbursed', 'sum')
    ).reset_index()
    g['conversion_pct'] = round(g['disbursals'] / g['leads'] * 100, 1)
    return g.sort_values('conversion_pct')

region_summary = funnel_summary(df, 'region')
segment_summary = funnel_summary(df, 'segment')
product_summary = funnel_summary(df, 'product')

#Root-cause flag
worst_region = region_summary.iloc[0]
worst_segment = segment_summary.iloc[0]
insight = (f"{worst_region['region']} shows the lowest conversion at "
           f"{worst_region['conversion_pct']}%, largely driven by underperformance "
           f"in the {worst_segment['segment']} segment ({worst_segment['conversion_pct']}% conversion).")

print(insight)

#Auto-generate PPT
prs = Presentation()

#Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Connector Channel Performance Review"
slide.placeholders[1].text = "Weekly MIS Report — Auto-generated"

#Slide 2: Region funnel chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Region-wise Conversion Funnel"
chart_data = CategoryChartData()
chart_data.categories = region_summary['region'].tolist()
chart_data.add_series('Conversion %', region_summary['conversion_pct'].tolist())
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(5), chart_data)

#Slide 3: Insight slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Insight — Root Cause"
slide.placeholders[1].text = insight

prs.save('Connector_Channel_Weekly_Report.pptx')
print("PPT report generated successfully.")